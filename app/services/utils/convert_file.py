import os
import tempfile
from pathlib import Path
from docx import Document
import pandas as pd
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import openpyxl
from pptx import Presentation
import csv

class FileConverter:
    def __init__(self):
        """Initialize file converter"""
        self.supported_formats = {
            '.docx': self.docx_to_pdf,
            '.doc': self.doc_to_pdf,
            '.xlsx': self.xlsx_to_pdf,
            '.xls': self.xls_to_pdf,
            '.csv': self.csv_to_pdf,
            '.txt': self.txt_to_pdf,
            '.pptx': self.pptx_to_pdf,
            '.ppt': self.ppt_to_pdf
        }
    
    def is_convertible(self, file_path):
        """Check if file can be converted to PDF"""
        extension = Path(file_path).suffix.lower()
        return extension in self.supported_formats
    
    def convert_to_pdf(self, input_path, output_path=None):
        """Convert file to PDF format"""
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
        
        extension = Path(input_path).suffix.lower()
        
        if extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {extension}")
        
        # Generate output path if not provided
        if output_path is None:
            input_stem = Path(input_path).stem
            output_path = os.path.join(tempfile.gettempdir(), f"{input_stem}_converted.pdf")
        
        print(f"FileConverter: Converting {extension} file to PDF...")
        
        # Call appropriate conversion function
        converter_func = self.supported_formats[extension]
        converter_func(input_path, output_path)
        
        print(f"FileConverter: Successfully converted to {os.path.basename(output_path)}")
        return output_path
    
    def docx_to_pdf(self, input_path, output_path):
        """Convert DOCX to PDF"""
        try:
            print(f"FileConverter: Processing DOCX file...")
            doc = Document(input_path)
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Process paragraphs
            paragraph_count = 0
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    p = Paragraph(paragraph.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 6))
                    paragraph_count += 1
            
            pdf_doc.build(story)
            print(f"FileConverter: Processed {paragraph_count} paragraphs from DOCX")
            
        except Exception as e:
            raise Exception(f"Error converting DOCX to PDF: {e}")
    
    def doc_to_pdf(self, input_path, output_path):
        """Convert DOC to PDF (requires python-docx2txt or similar)"""
        try:
            print(f"FileConverter: Processing DOC file...")
            # For .doc files, we need additional libraries like python-docx2txt
            # This is a simplified version - you might need to install additional packages
            import docx2txt
            text = docx2txt.process(input_path)
            
            # Create PDF from extracted text
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            paragraphs = text.split('\n')
            paragraph_count = 0
            for para in paragraphs:
                if para.strip():
                    p = Paragraph(para, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 6))
                    paragraph_count += 1
            
            pdf_doc.build(story)
            print(f"FileConverter: Processed {paragraph_count} paragraphs from DOC")
            
        except ImportError:
            raise Exception("python-docx2txt package required for .doc files. Install with: pip install docx2txt")
        except Exception as e:
            raise Exception(f"Error converting DOC to PDF: {e}")
    
    def xlsx_to_pdf(self, input_path, output_path):
        """Convert XLSX to PDF"""
        try:
            print(f"FileConverter: Processing XLSX file...")
            # Read Excel file
            df = pd.read_excel(input_path, sheet_name=None)  # Read all sheets
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            sheet_count = 0
            total_rows = 0
            
            for sheet_name, sheet_df in df.items():
                # Add sheet title
                title = Paragraph(f"Sheet: {sheet_name}", styles['Heading1'])
                story.append(title)
                story.append(Spacer(1, 12))
                
                # Convert dataframe to table
                data = [sheet_df.columns.tolist()] + sheet_df.values.tolist()
                
                # Create table with limited width to fit page
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('FONTSIZE', (0, 1), (-1, -1), 6),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(table)
                story.append(Spacer(1, 20))
                
                sheet_count += 1
                total_rows += len(sheet_df)
            
            pdf_doc.build(story)
            print(f"FileConverter: Processed {sheet_count} sheets with {total_rows} total rows from XLSX")
            
        except Exception as e:
            raise Exception(f"Error converting XLSX to PDF: {e}")
    
    def xls_to_pdf(self, input_path, output_path):
        """Convert XLS to PDF"""
        print(f"FileConverter: Processing XLS file (using XLSX converter)...")
        # Similar to xlsx_to_pdf but for older Excel format
        self.xlsx_to_pdf(input_path, output_path)
    
    def csv_to_pdf(self, input_path, output_path):
        """Convert CSV to PDF"""
        try:
            print(f"FileConverter: Processing CSV file...")
            # Read CSV file
            df = pd.read_csv(input_path)
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title
            title = Paragraph("CSV Data", styles['Heading1'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            # Convert dataframe to table
            data = [df.columns.tolist()] + df.values.tolist()
            
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 8),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('FONTSIZE', (0, 1), (-1, -1), 6),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(table)
            pdf_doc.build(story)
            
            print(f"FileConverter: Processed {len(df)} rows from CSV")
            
        except Exception as e:
            raise Exception(f"Error converting CSV to PDF: {e}")
    
    def txt_to_pdf(self, input_path, output_path):
        """Convert TXT to PDF"""
        try:
            print(f"FileConverter: Processing TXT file...")
            with open(input_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            paragraphs = text.split('\n')
            paragraph_count = 0
            for para in paragraphs:
                if para.strip():
                    p = Paragraph(para, styles['Normal'])
                    story.append(p)
                    paragraph_count += 1
                story.append(Spacer(1, 6))
            
            pdf_doc.build(story)
            print(f"FileConverter: Processed {paragraph_count} paragraphs from TXT")
            
        except Exception as e:
            raise Exception(f"Error converting TXT to PDF: {e}")
    
    def pptx_to_pdf(self, input_path, output_path):
        """Convert PPTX to PDF"""
        try:
            print(f"FileConverter: Processing PPTX file...")
            prs = Presentation(input_path)
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            slide_count = 0
            text_count = 0
            
            for i, slide in enumerate(prs.slides):
                # Add slide number
                slide_title = Paragraph(f"Slide {i + 1}", styles['Heading1'])
                story.append(slide_title)
                story.append(Spacer(1, 12))
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        p = Paragraph(shape.text, styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                        text_count += 1
                
                story.append(Spacer(1, 20))
                slide_count += 1
            
            pdf_doc.build(story)
            print(f"FileConverter: Processed {slide_count} slides with {text_count} text elements from PPTX")
            
        except Exception as e:
            raise Exception(f"Error converting PPTX to PDF: {e}")
    
    def ppt_to_pdf(self, input_path, output_path):
        """Convert PPT to PDF (requires additional libraries)"""
        try:
            print(f"FileConverter: Processing PPT file...")
            # For .ppt files, you might need win32com (Windows only) or other libraries
            # This is a placeholder - actual implementation depends on your system
            raise Exception("PPT conversion requires additional setup. Consider converting to PPTX first.")
        except Exception as e:
            raise Exception(f"Error converting PPT to PDF: {e}")


def convert_file_to_pdf(file_path, output_path=None):
    """Main function to convert file to PDF"""
    converter = FileConverter()
    
    if not converter.is_convertible(file_path):
        extension = Path(file_path).suffix.lower()
        raise ValueError(f"File format {extension} is not supported for conversion")
    
    try:
        pdf_path = converter.convert_to_pdf(file_path, output_path)
        print(f"Successfully converted {file_path} to {pdf_path}")
        return pdf_path
    except Exception as e:
        print(f"Conversion failed: {e}")
        raise e
